import os
import subprocess
import tempfile

import docx
import fitz  # PyMuPDF
import pptx

from app.services.ocr.ocr import run_tesseract_ocr


def extract_text_from_txt(file_path: str) -> str:
    """Extract raw string text from text files."""
    try:
        with open(file_path, encoding="utf-8", errors="ignore") as f:
            return f.read().strip()
    except Exception as e:
        print(f"Error reading TXT file {file_path}: {e}")
        return ""


def extract_text_from_docx(file_path: str) -> str:
    """Extract raw text paragraphs from Microsoft Word documents."""
    try:
        doc = docx.Document(file_path)
        paragraphs = [p.text for p in doc.paragraphs]
        return "\n".join(paragraphs).strip()
    except Exception as e:
        print(f"Error reading DOCX file {file_path}: {e}")
        return ""


def _convert_ppt_to_pptx(file_path: str) -> str | None:
    """Convert old .ppt to .pptx using LibreOffice headless mode. Returns path to converted file or None."""
    try:
        out_dir = tempfile.mkdtemp()
        print(f"Starting LibreOffice conversion: {file_path} -> {out_dir}")
        result = subprocess.run(
            [
                "soffice",
                "--headless",
                "--convert-to",
                "pptx",
                "--outdir",
                out_dir,
                file_path,
            ],
            capture_output=True,
            text=True,
            timeout=180,
        )
        print(f"LibreOffice exit code: {result.returncode}, stderr: {result.stderr[:200]}")
        if result.returncode != 0:
            print(f"LibreOffice conversion failed: {result.stderr}")
            return None
        base = os.path.splitext(os.path.basename(file_path))[0]
        converted = os.path.join(out_dir, f"{base}.pptx")
        exists = os.path.exists(converted)
        print(f"Converted file exists: {exists}, path: {converted}")
        return converted if exists else None
    except subprocess.TimeoutExpired:
        print(f"LibreOffice conversion timed out after 180s for {file_path}")
        return None
    except Exception as e:
        print(f"Error converting PPT to PPTX: {e}")
        return None


def extract_text_from_pptx(file_path: str) -> str:
    """Extract text lines from Microsoft PowerPoint presentations."""
    target = file_path
    needs_cleanup = False
    try:
        prs = pptx.Presentation(target)
    except Exception:
        # Fallback: try converting old .ppt via LibreOffice
        print(f"python-pptx failed on {target}, attempting LibreOffice conversion...")
        converted = _convert_ppt_to_pptx(target)
        if converted is None:
            print(f"Could not convert PPT to PPTX for {target}")
            return ""
        target = converted
        needs_cleanup = True
        try:
            prs = pptx.Presentation(target)
        except Exception as e:
            print(f"Error reading converted PPTX file {target}: {e}")
            return ""

    try:
        text_runs = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text_runs.append(shape.text)
        return "\n".join(text_runs).strip()
    except Exception as e:
        print(f"Error reading PPTX file {target}: {e}")
        return ""
    finally:
        if needs_cleanup and target != file_path:
            try:
                os.remove(target)
                os.rmdir(os.path.dirname(target))
            except Exception:
                pass


def extract_text_from_pdf(file_path: str, ocr_enabled: bool = True) -> str:
    """Extract text from PDF using PyMuPDF. Runs OCR if text content is low."""
    try:
        doc = fitz.open(file_path)
        text_runs = []

        # 1. Standard text extraction
        for page in doc:
            text_runs.append(page.get_text())

        full_text = "\n".join(text_runs).strip()

        # 2. Check if we need to fall back to OCR
        if len(full_text) < 150 and ocr_enabled:
            print(f"Low text density detected in PDF ({len(full_text)} chars). Running Tesseract OCR fallback...")
            ocr_text_runs = []

            # Limit to first 10 pages to prevent blocking CPU indefinitely
            pages_to_ocr = min(len(doc), 10)
            for page_num in range(pages_to_ocr):
                page = doc.load_page(page_num)
                # Render page to image pixmap
                pix = page.get_pixmap(dpi=150)
                temp_img_path = os.path.join(
                    tempfile.gettempdir(),
                    f"page_{page_num}_{os.path.basename(file_path)}.png",
                )
                pix.save(temp_img_path)

                # Perform OCR
                page_text = run_tesseract_ocr(temp_img_path)
                ocr_text_runs.append(page_text)

                # Clean up temporary page image
                if os.path.exists(temp_img_path):
                    os.remove(temp_img_path)

            full_text = "\n".join(ocr_text_runs).strip()

        return full_text

    except Exception as e:
        print(f"Error reading PDF file {file_path}: {e}")
        return ""


def extract_content(file_path: str, file_type: str, ocr_enabled: bool = True) -> str:
    """Routes the file to the appropriate format text extractor."""
    ext = file_type.upper()
    if ext == "PDF":
        return extract_text_from_pdf(file_path, ocr_enabled)
    elif ext == "DOCX":
        return extract_text_from_docx(file_path)
    elif ext in ("PPT", "PPTX"):
        return extract_text_from_pptx(file_path)
    elif ext == "TXT":
        return extract_text_from_txt(file_path)
    elif ext in ["PNG", "JPG", "JPEG"]:
        return run_tesseract_ocr(file_path)
    else:
        return ""
